import os
import pypdf
from pptx import Presentation
from pathlib import Path
import logging
from supabase import create_client
from langchain_openai import OpenAIEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from app.core.config import settings

logger = logging.getLogger(__name__)

UPLOAD_DIR = Path("data/uploads")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

class RAGService:
    def __init__(self):
        self.client = None
        self.embeddings = None

        if settings.SUPABASE_URL and settings.SUPABASE_SERVICE_KEY:
            self.client = create_client(settings.SUPABASE_URL, settings.SUPABASE_SERVICE_KEY)
            logger.info("RAG Service connected to Supabase")
        else:
            logger.warning("Supabase credentials not set. RAG service database sync disabled.")

        api_key = settings.OPENROUTER_API_KEY
        if api_key:
            self.embeddings = OpenAIEmbeddings(
                api_key=api_key,
                model="openai/text-embedding-3-small",
                base_url="https://openrouter.ai/api/v1"
            )
        else:
            logger.warning("OPENROUTER_API_KEY not set. Falling back or failing embedding generation.")

        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )

    def _extract_text(self, filepath: Path) -> str:
        try:
            if filepath.suffix.lower() == ".pdf":
                text = ""
                with open(filepath, "rb") as f:
                    pdf = pypdf.PdfReader(f)
                    for page in pdf.pages:
                        extracted = page.extract_text()
                        if extracted:
                            text += extracted + "\n"
                return text
            elif filepath.suffix.lower() == ".pptx":
                text = ""
                prs = Presentation(filepath)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text += shape.text + "\n"
                return text
            else:
                with open(filepath, "r", encoding="utf-8") as f:
                    return f.read()
        except Exception as e:
            logger.error(f"Error extracting text from {filepath.name}: {e}")
            return ""

    def add_document(self, filepath: Path):
        if not self.client or not self.embeddings:
            logger.error("RAG Service not initialized properly (Missing Supabase/Embedding config).")
            return

        content = self._extract_text(filepath)
        if not content.strip():
            logger.warning(f"No text extracted from {filepath.name}")
            return

        chunks = self.text_splitter.split_text(content)
        
        filename = filepath.name
        total_chunks = len(chunks)
        logger.info(f"Extracting {total_chunks} chunks for {filename}. Batching embedding generation...")

        batch_size = 100
        for i in range(0, total_chunks, batch_size):
            batch = chunks[i:i+batch_size]
            try:
                # Process in batches to massively speed up 1000 page PDFs like textbooks
                vectors = self.embeddings.embed_documents(batch)
                
                data = [
                    {
                        "content": b,
                        "embedding": v,
                        "metadata": {"source": filename},
                    }
                    for b, v in zip(batch, vectors)
                ]
                
                self.client.table("documents").insert(data).execute()
                logger.info(f"Supabase Vector Indexed batch {i//batch_size + 1}/{(total_chunks-1)//batch_size + 1} for {filename}.")
            except Exception as e:
                logger.error(f"Failed to index batch {i//batch_size + 1} for {filename}: {e}")

        logger.info(f"Finished indexing all {total_chunks} chunks for {filename}.")

    def search(self, query: str, limit: int = 3) -> list[str]:
        if not self.client or not self.embeddings:
            return []
            
        try:
            vector = self.embeddings.embed_query(query)

            result = self.client.rpc("match_documents", {
                "query_embedding": vector,
                "match_count": limit,
            }).execute()

            docs = []
            for row in (result.data or []):
                source = row.get("metadata", {}).get("source", "Unknown")
                docs.append(f"[From {source}]:\n{row['content']}")
                
            return docs
        except Exception as e:
            logger.error(f"RAG search error: {e}")
            return []

rag_service = RAGService()
